#!/usr/bin/env python3
"""Generate partner Zoom deck: PowerPoint for Leads.Konversus."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from copy import deepcopy

OUT = "/var/www/www-root/data/www/leads.konversus.ru/public/presentations/partner-zoom/Leads-Konversus-Partner-Zoom.pptx"

# Brand
NAVY = RGBColor(0x0F, 0x1C, 0x2E)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DK = RGBColor(0x0F, 0x76, 0x6E)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED_SOFT = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)


def set_run(run, size=18, bold=False, color=NAVY, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_text(shape, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def rounded(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def footer(slide, page, total=13):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"leads.konversus.ru  ·  {page}/{total}"
    set_run(run, size=11, color=SLATE)


def title_bar(slide, title, subtitle=None):
    fill_rect(slide, 0, 0, Inches(13.333), Inches(1.15), NAVY)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.55))
    add_text(box, title, size=28, bold=True, color=WHITE)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.35))
        add_text(box2, subtitle, size=14, color=RGBColor(0x94, 0xA3, 0xB8))


def bullet_card(slide, left, top, w, h, title, lines, accent=TEAL):
    card = rounded(slide, left, top, w, h, CREAM)
    bar = fill_rect(slide, left, top, Inches(0.12), h, accent)
    t = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), w - Inches(0.45), Inches(0.4))
    add_text(t, title, size=16, bold=True, color=NAVY)
    body = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55), w - Inches(0.45), h - Inches(0.7))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        run = p.add_run()
        run.text = "•  " + line
        set_run(run, size=13, color=SLATE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # —— 1 Title ——
    s = blank_slide(prs)
    fill_rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY)
    fill_rect(s, 0, Inches(6.6), Inches(13.333), Inches(0.9), TEAL)
    logo = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(6), Inches(0.4))
    add_text(logo, "◈  Leads.Konversus", size=18, bold=True, color=TEAL)
    t = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(7.8), Inches(2.2))
    tf = add_text(
        t,
        "Все заказы с Profi — только ваши,\nтолько квалифицированные",
        size=32,
        bold=True,
        color=WHITE,
    )
    add_para(tf, "ВЫ ВЫБИРАЕТЕ ЗАДАЧУ И КЛИЕНТА, ЗА КОТОРОГО ПЛАТИТЕ", size=16, bold=True, color=TEAL, space_before=14)
    sub = s.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(7.5), Inches(1.0))
    add_text(sub, "25+ реальных клиентов в день,\nкоторым можно продать и заработать", size=20, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Mock TG card
    card = rounded(s, Inches(8.6), Inches(1.4), Inches(4.2), Inches(4.6), WHITE)
    th = s.shapes.add_textbox(Inches(8.85), Inches(1.55), Inches(3.7), Inches(0.35))
    add_text(th, "Telegram · заявка", size=12, bold=True, color=TEAL)
    lines = [
        ("══ PROFI.RU ══", NAVY, True, 13),
        ("Создание сайта на Тильде", NAVY, True, 15),
        ("💰 до 40 000 ₽", TEAL_DK, True, 14),
        ("📍 Москва · ⭐ 12 отз.", SLATE, False, 12),
        ("Клиент изучает цены", SLATE, False, 12),
        ("", SLATE, False, 10),
        ("[ Открыть заказ ]", ORANGE, True, 13),
    ]
    box = s.shapes.add_textbox(Inches(8.85), Inches(2.1), Inches(3.7), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, col, bold, sz in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = text
        set_run(run, size=sz, bold=bold, color=col)
    cta = s.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(8), Inches(0.4))
    add_text(cta, "Презентация для Zoom  ·  leads.konversus.ru", size=14, bold=True, color=WHITE)

    # —— 2A Builders ——
    s = blank_slide(prs)
    title_bar(s, "Кому мы даём поток заказов — строителям и ремонтникам", "Вариант А")
    bullet_card(
        s,
        Inches(0.4),
        Inches(1.5),
        Inches(6.1),
        Inches(2.4),
        "Частные мастера",
        [
            "Сантехники, электрики, плиточники, штукатуры, маляры…",
            "Заказы строго по специализации — без мусора",
        ],
        TEAL,
    )
    bullet_card(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(6.1),
        Inches(2.4),
        "Бригады и подрядчики",
        [
            "От 2–3 человек до крупных команд",
            "Поток объектов; фильтр: жилой/коммерция, бюджет, сроки",
        ],
        TEAL_DK,
    )
    who = rounded(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.2), CREAM)
    t = s.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(12), Inches(1.8))
    tf = add_text(t, "Кто ищет на Profi", size=16, bold=True, color=NAVY)
    add_para(
        tf,
        "Ремонт квартир и домов · отделка · строительство · сантехника · электрика · фасады · ландшафт и т.д.",
        size=15,
        color=SLATE,
        space_before=10,
    )
    footer(s, 2)

    # —— 2B Legal ——
    s = blank_slide(prs)
    title_bar(s, "Кому мы даём поток клиентов — юристам и бухгалтерам", "Вариант Б")
    bullet_card(
        s,
        Inches(0.4),
        Inches(1.5),
        Inches(6.1),
        Inches(2.4),
        "Частные практики (ИП, самозанятые)",
        [
            "Консультации, договоры, декларации",
            "Только под вашу специализацию",
        ],
        TEAL,
    )
    bullet_card(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(6.1),
        Inches(2.4),
        "Юр. и бухгалтерские фирмы",
        [
            "Поток корпоративных клиентов",
            "Фильтр: бюджет, услуга, срочность, рейтинг заказчика",
        ],
        TEAL_DK,
    )
    who = rounded(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.2), CREAM)
    t = s.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(12), Inches(1.8))
    tf = add_text(t, "Кто ищет на Profi", size=16, bold=True, color=NAVY)
    add_para(
        tf,
        "Юридические консультации · бухобслуживание · аудит · налоги · регистрация ООО · суд и т.д.",
        size=15,
        color=SLATE,
        space_before=10,
    )
    footer(s, 3)

    # —— 3 Problem ——
    s = blank_slide(prs)
    title_bar(s, "Что вы теряете на Profi каждый день")
    bullets = [
        "Profi показывает 2–3 заявки в день — искусственное ограничение алгоритма",
        "Остальные ~90% заказов вы не видите — конкуренты успевают раньше",
        "Без фильтрации: низкий бюджет, не ваш тип работ, боты, «холодные» клиенты",
    ]
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.2), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(14)
        run = p.add_run()
        run.text = "▸  " + b
        set_run(run, size=16, color=NAVY)

    # Big number
    big = rounded(s, Inches(8.0), Inches(1.5), Inches(4.8), Inches(2.0), NAVY)
    t = s.shapes.add_textbox(Inches(8.2), Inches(1.7), Inches(4.4), Inches(1.6))
    tf = add_text(t, "Вы видите лишь", size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
    add_para(tf, "10%", size=48, bold=True, color=ORANGE, space_before=4)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, "от всех заказов", size=14, color=WHITE, space_before=2)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    # Pie-like chart: 10% vs 90%
    chart_data = CategoryChartData()
    chart_data.categories = ["Видите (10%)", "Скрыто (90%)"]
    chart_data.add_series("Заказы", (10, 90))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(8.0), Inches(3.8), Inches(4.8), Inches(2.8), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    footer(s, 4)

    # —— 4 Solution ——
    s = blank_slide(prs)
    title_bar(s, "Мы даём вам все заказы — и только те, что подходят", "Leads.Konversus")
    points = [
        ("20–25+", "заявок в день, а не 2–3"),
        ("15+", "параметров фильтрации под ваш бизнес"),
        ("сек.", "доставка в Telegram — вы первыми"),
        ("100–400 ₽", "цена квалифицированной заявки"),
    ]
    x = 0.4
    for num, label in points:
        card = rounded(s, Inches(x), Inches(1.5), Inches(3.0), Inches(1.9), CREAM)
        t = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.7), Inches(2.7), Inches(1.5))
        tf = add_text(t, num, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_para(tf, label, size=13, color=SLATE, space_before=8)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        x += 3.2

    # Flow boxes
    steps = [("Profi", "Лента заказов"), ("Сервис", "Сбор + защита"), ("Фильтры", "15+ правил"), ("Telegram", "Карточка вам")]
    x = 0.5
    for i, (a, b) in enumerate(steps):
        card = rounded(s, Inches(x), Inches(4.0), Inches(2.6), Inches(1.8), NAVY if i % 2 == 0 else TEAL)
        t = s.shapes.add_textbox(Inches(x + 0.1), Inches(4.3), Inches(2.4), Inches(1.3))
        tf = add_text(t, a, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_para(tf, b, size=13, color=RGBColor(0xE2, 0xE8, 0xF0), space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        if i < 3:
            arr = s.shapes.add_textbox(Inches(x + 2.55), Inches(4.55), Inches(0.4), Inches(0.5))
            add_text(arr, "→", size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        x += 3.15
    footer(s, 5)

    # —— 5 Filters ——
    s = blank_slide(prs)
    title_bar(s, "Никаких сложных настроек — мы всё сделаем за вас")
    left = [
        "Анализируем ваш профиль и типичные заказы",
        "Бюджет, география, виды работ — под вас",
        "Ключевые слова, рейтинг, тип объекта, сроки…",
        "Вы просто получаете поток в Telegram",
    ]
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in left:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(16)
        run = p.add_run()
        run.text = "✓  " + line
        set_run(run, size=17, color=NAVY)

    mock = rounded(s, Inches(7.4), Inches(1.5), Inches(5.4), Inches(4.8), WHITE)
    # border effect via navy top
    fill_rect(s, Inches(7.4), Inches(1.5), Inches(5.4), Inches(0.5), TEAL)
    ht = s.shapes.add_textbox(Inches(7.6), Inches(1.58), Inches(5), Inches(0.35))
    add_text(ht, "Пример карточки после фильтров", size=13, bold=True, color=WHITE)
    body = s.shapes.add_textbox(Inches(7.7), Inches(2.2), Inches(4.9), Inches(3.8))
    tf = body.text_frame
    tf.word_wrap = True
    rows = [
        ("Тип работ", "Электрика · квартира", True),
        ("Бюджет", "от 15 000 ₽", True),
        ("Рейтинг", "⭐ клиент с отзывами", True),
        ("Город", "Москва / МО", False),
        ("Минус", "отсеян «дешёвый» мусор", False),
    ]
    first = True
    for k, v, hi in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = k + ": "
        set_run(r1, size=13, bold=True, color=SLATE)
        r2 = p.add_run()
        r2.text = v
        set_run(r2, size=14, bold=hi, color=TEAL_DK if hi else NAVY)
    footer(s, 6)

    # —— 6 Security ——
    s = blank_slide(prs)
    title_bar(s, "Работаем без риска бана", "Подстраховка аккаунта Profi")
    note = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.0))
    tf = add_text(
        note,
        "Profi банит грубых парсеров. Поэтому: эмуляция человека, случайные паузы, отдельный контур на партнёра.",
        size=16,
        color=NAVY,
    )
    levels = [
        ("Light", "Базовая защита", "Тариф «Мастер»", TEAL),
        ("Balanced", "Усиленный режим", "Тариф «Команда»", TEAL_DK),
        ("Stealth", "Максимальная осторожность", "Тариф «Бизнес»", NAVY),
    ]
    x = 0.5
    for name, desc, tar, col in levels:
        card = rounded(s, Inches(x), Inches(2.7), Inches(3.9), Inches(3.2), CREAM)
        fill_rect(s, Inches(x), Inches(2.7), Inches(3.9), Inches(0.7), col)
        t = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.85), Inches(3.5), Inches(0.45))
        add_text(t, "🛡  " + name, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        b = s.shapes.add_textbox(Inches(x + 0.25), Inches(3.7), Inches(3.4), Inches(1.8))
        tf = add_text(b, desc, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_para(tf, tar, size=14, color=SLATE, space_before=12)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, "Имитация живого пользователя", size=12, color=SLATE, space_before=10)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        x += 4.2
    footer(s, 7)

    # —— 7 Results ——
    s = blank_slide(prs)
    title_bar(s, "Что меняется после подключения")
    chart_data = CategoryChartData()
    chart_data.categories = ["Было (Profi)", "Стало (Leads)"]
    chart_data.add_series("Заявок в день", (2.5, 25))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.4), Inches(6.5), Inches(5.0), chart_data
    ).chart
    chart.has_legend = False

    stats = [
        ("25+", "проверенных лидов в день"),
        ("100–400 ₽", "цена заявки"),
        ("время", "не скроллите ленту"),
        ("конверсия", "только платёжеспособные"),
    ]
    y = 1.5
    for a, b in stats:
        card = rounded(s, Inches(7.3), Inches(y), Inches(5.5), Inches(1.15), CREAM)
        t = s.shapes.add_textbox(Inches(7.5), Inches(y + 0.15), Inches(5.1), Inches(0.9))
        tf = add_text(t, a, size=22, bold=True, color=TEAL)
        add_para(tf, b, size=14, color=SLATE, space_before=2)
        y += 1.3
    footer(s, 8)

    # —— 8 Comparison table ——
    s = blank_slide(prs)
    title_bar(s, "Leads.Konversus vs другие способы")
    rows = [
        ["", "Обычный Profi", "Контекстная реклама", "Leads.Konversus"],
        ["Заявок в день", "2–3", "Зависит от бюджета", "20–25+"],
        ["Фильтрация", "Нет", "Частичная", "15+ параметров"],
        ["Цена за заявку", "«Бесплатно», но нет потока", "500–1500 ₽", "100–400 ₽"],
        ["Риск бана", "Нет", "Нет", "Минимальный (антидетект)"],
    ]
    table = s.shapes.add_table(len(rows), 4, Inches(0.4), Inches(1.6), Inches(12.5), Inches(4.6)).table
    table.columns[0].width = Inches(2.4)
    for c in range(1, 4):
        table.columns[c].width = Inches(3.366)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(
                        run,
                        size=13 if r else 12,
                        bold=(r == 0 or c == 0 or c == 3),
                        color=WHITE if r == 0 or c == 3 else NAVY,
                    )
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            # fill
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = NAVY
            elif c == 3:
                fill.fore_color.rgb = RGBColor(0xCC, 0xFB, 0xF1)  # teal soft
            elif r % 2 == 0:
                fill.fore_color.rgb = CREAM
            else:
                fill.fore_color.rgb = WHITE
    footer(s, 9)

    # —— 9 Tariffs ——
    s = blank_slide(prs)
    title_bar(s, "Выберите свой тариф")
    tariffs = [
        ("Мастер", "19 990 ₽/мес", "Частники", ["До 5 фильтров", "1 Telegram", "Антидетект Light"], TEAL),
        ("Команда", "30 000 ₽/мес", "Малые организации", ["До 10 фильтров", "2 менеджера в TG", "Антидетект Balanced"], TEAL_DK),
        ("Бизнес", "50 000 ₽/мес", "Крупные организации", ["Фильтры без лимита", "Любое число менеджеров", "API + Stealth"], NAVY),
    ]
    x = 0.4
    for name, price, who, feats, col in tariffs:
        card = rounded(s, Inches(x), Inches(1.4), Inches(4.0), Inches(4.4), CREAM)
        fill_rect(s, Inches(x), Inches(1.4), Inches(4.0), Inches(1.35), col)
        t = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.5), Inches(3.6), Inches(1.15))
        tf = add_text(t, name, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_para(tf, price, size=18, bold=True, color=WHITE, space_before=4)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, who, size=12, color=RGBColor(0xE2, 0xE8, 0xF0), space_before=2)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        b = s.shapes.add_textbox(Inches(x + 0.3), Inches(3.0), Inches(3.5), Inches(2.2))
        tf = b.text_frame
        first = True
        for f in feats:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(10)
            run = p.add_run()
            run.text = "•  " + f
            set_run(run, size=14, color=NAVY)
        x += 4.25
    trial = rounded(s, Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.7), ORANGE)
    t = s.shapes.add_textbox(Inches(0.6), Inches(6.18), Inches(12.1), Inches(0.45))
    add_text(
        t,
        "Пробный период: 7 дней за 7 000 ₽  ·  без привязки к основному тарифу",
        size=16,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 10)

    # —— 10 Six reasons ——
    s = blank_slide(prs)
    title_bar(s, "6 причин выбрать Leads.Konversus")
    reasons = [
        ("Полнота", "Все заказы, не малая часть"),
        ("Точность", "Фильтры под ваш бизнес"),
        ("Скорость", "Telegram — вы первые"),
        ("Экономия", "100–400 ₽ vs 500–1500 ₽"),
        ("Безопасность", "Аккаунт Profi под защитой"),
        ("Простота", "Мы настраиваем за вас"),
    ]
    positions = [
        (0.4, 1.5),
        (4.7, 1.5),
        (9.0, 1.5),
        (0.4, 4.0),
        (4.7, 4.0),
        (9.0, 4.0),
    ]
    for (title, desc), (lx, ty) in zip(reasons, positions):
        card = rounded(s, Inches(lx), Inches(ty), Inches(4.0), Inches(2.1), CREAM)
        fill_rect(s, Inches(lx), Inches(ty), Inches(0.15), Inches(2.1), TEAL)
        t = s.shapes.add_textbox(Inches(lx + 0.35), Inches(ty + 0.4), Inches(3.4), Inches(1.4))
        tf = add_text(t, title, size=20, bold=True, color=NAVY)
        add_para(tf, desc, size=14, color=SLATE, space_before=8)
    footer(s, 11)

    # —— 11 CTA ——
    s = blank_slide(prs)
    fill_rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.5))
    add_text(
        t,
        "Начните получать 25+ квалифицированных\nзаявок уже сегодня",
        size=30,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    sub = s.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.2))
    tf = add_text(
        sub,
        "Пробная неделя за 7 000 ₽ — убедитесь, что заявки есть и они качественные.",
        size=18,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )
    add_para(tf, "Сайт: leads.konversus.ru  ·  напишите нам в Telegram", size=16, color=TEAL, space_before=12)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    btn = rounded(s, Inches(4.2), Inches(5.0), Inches(4.9), Inches(0.9), ORANGE)
    bt = s.shapes.add_textbox(Inches(4.2), Inches(5.2), Inches(4.9), Inches(0.55))
    add_text(bt, "Попробовать неделю →", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # —— 12 Contacts ——
    s = blank_slide(prs)
    fill_rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY)
    fill_rect(s, 0, Inches(5.8), Inches(13.333), Inches(1.7), TEAL)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0))
    add_text(t, "Остались вопросы? Мы на связи", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    c = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.8))
    tf = add_text(c, "◈  Leads.Konversus", size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_para(tf, "Сайт:  leads.konversus.ru", size=18, color=WHITE, space_before=16)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, "Telegram:  напишите менеджеру / боту сервиса", size=16, color=RGBColor(0xCB, 0xD5, 0xE1), space_before=8)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    f = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.8))
    add_text(
        f,
        "Спасибо за внимание  ·  Zoom-презентация для партнёров",
        size=16,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    prs.save(OUT)
    print("Wrote", OUT)


if __name__ == "__main__":
    build()
